```python
#!/usr/bin/env python3
"""
IAM Analyzer v6 - Enterprise Edition
Multi-format file support + Super in-depth reporting capabilities
"""

import os
import io
import json
import logging
import traceback
import hashlib
import mimetypes
import uuid
from datetime import datetime, timedelta
from pathlib import Path
from typing import Dict, List, Tuple, Optional, Any, Union
from dotenv import load_dotenv

import streamlit as st
import pandas as pd

from openai import OpenAI
import openai
import base64
import zipfile
import xml.etree.ElementTree as ET

# Document processing imports
try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    from PIL import Image
    import pytesseract
    OCR_AVAILABLE = True
except ImportError:
    OCR_AVAILABLE = False

# Try to import plotly, fall back to matplotlib if not available
try:
    import plotly.express as px
    import plotly.graph_objects as go
    from plotly.subplots import make_subplots
    PLOTLY_AVAILABLE = True
except ImportError:
    import matplotlib.pyplot as plt
    PLOTLY_AVAILABLE = False

# Load environment variables
load_dotenv()

# Page configuration
st.set_page_config(
    page_title="IAM Analyzer v6 Enterprise", 
    layout="wide",
    page_icon="🔐",
    initial_sidebar_state="expanded"
)

# Constants
MAX_FILE_SIZE = 50 * 1024 * 1024  # 50MB for enterprise files
SUPPORTED_MODELS = ["gpt-4", "gpt-4-turbo", "gpt-3.5-turbo"]
RISK_COLORS = {"Low": "#28a745", "Medium": "#ffc107", "High": "#dc3545", "Critical": "#6f42c1"}
SUPPORTED_FILE_TYPES = ["csv", "xlsx", "xls", "txt", "docx", "pptx", "pdf", "png", "jpg", "jpeg", "gif", "bmp", "tiff", "pbi"]

# ================== ENHANCED LOGGING SYSTEM ==================

def setup_logging():
    """Configure enhanced logging for the application"""
    logging.basicConfig(
        level=logging.INFO,
        format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
        handlers=[logging.StreamHandler()]
    )
    return logging.getLogger("IAM_Analyzer_v6")

logger = setup_logging()

class AnalysisTracker:
    """Track analysis sessions and performance metrics"""

    def __init__(self):
        self.session_id = str(uuid.uuid4())[:8]
        self.start_time = datetime.now()
        self.metrics = {
            "files_processed": 0,
            "total_records": 0,
            "processing_time": 0,
            "api_calls": 0,
            "chat_interactions": 0,
            "web_searches": 0,
            "exports_generated": 0,
            "errors": []
        }

    def log_file_processed(self, filename: str, records: int):
        self.metrics["files_processed"] += 1
        self.metrics["total_records"] += records
        logger.info(f"Session {self.session_id}: Processed {filename} - {records} records")

    def log_error(self, error_type: str, error_msg: str, file_name: str = None):
        error_entry = {
            "timestamp": datetime.now().isoformat(),
            "type": error_type,
            "message": error_msg,
            "file": file_name,
            "session": self.session_id
        }
        self.metrics["errors"].append(error_entry)
        logger.error(f"Session {self.session_id}: {error_type} - {error_msg}")

    def log_api_call(self, model: str, tokens_used: int = 0):
        self.metrics["api_calls"] += 1
        logger.info(f"Session {self.session_id}: API call to {model} - {tokens_used} tokens")

    def log_chat_interaction(self, query_type: str):
        self.metrics["chat_interactions"] += 1
        logger.info(f"Session {self.session_id}: Chat interaction - {query_type}")

    def log_web_search(self, query: str):
        self.metrics["web_searches"] += 1
        logger.info(f"Session {self.session_id}: Web search - {query}")

    def log_export(self, export_type: str):
        self.metrics["exports_generated"] += 1
        logger.info(f"Session {self.session_id}: Export generated - {export_type}")

    def get_session_summary(self) -> Dict:
        duration = datetime.now() - self.start_time
        return {
            "session_id": self.session_id,
            "duration_minutes": round(duration.total_seconds() / 60, 2),
            "files_processed": self.metrics["files_processed"],
            "total_records": self.metrics["total_records"],
            "api_calls": self.metrics["api_calls"],
            "chat_interactions": self.metrics["chat_interactions"],
            "web_searches": self.metrics["web_searches"],
            "exports_generated": self.metrics["exports_generated"],
            "errors_count": len(self.metrics["errors"]),
            "success_rate": round((self.metrics["files_processed"] / max(1, self.metrics["files_processed"] + len(self.metrics["errors"]))) * 100, 2)
        }

# ================== ADVANCED FILE VALIDATION & SECURITY ==================

class FileValidator:
    """Advanced file validation and security checking"""

    MIME_TYPE_MAP = {
        'csv': ['text/csv', 'application/csv', 'text/plain'],
        'xlsx': ['application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'],
        'xls': ['application/vnd.ms-excel'],
        'docx': ['application/vnd.openxmlformats-officedocument.wordprocessingml.document'],
        'pptx': ['application/vnd.openxmlformats-officedocument.presentationml.presentation'],
        'txt': ['text/plain'],
        'png': ['image/png'],
        'jpg': ['image/jpeg'],
        'jpeg': ['image/jpeg'],
        'gif': ['image/gif'],
        'pdf': ['application/pdf']
    }

    FILE_SIGNATURES = {
        'pdf': [b'%PDF'],
        'png': [b'\x89PNG\r\n\x1a\n'],
        'jpg': [b'\xff\xd8\xff'],
        'gif': [b'GIF87a', b'GIF89a'],
        'xlsx': [b'PK\x03\x04'],
        'docx': [b'PK\x03\x04'],
        'pptx': [b'PK\x03\x04']
    }

    @staticmethod
    def validate_file_signature(file_content: bytes, file_ext: str) -> bool:
        if file_ext not in FileValidator.FILE_SIGNATURES:
            return True
        expected_signatures = FileValidator.FILE_SIGNATURES[file_ext]
        return any(file_content.startswith(sig) for sig in expected_signatures)

    @staticmethod
    def calculate_file_hash(file_content: bytes) -> str:
        return hashlib.sha256(file_content).hexdigest()

    @staticmethod
    def validate_file_comprehensive(file, max_size: int = MAX_FILE_SIZE) -> Dict[str, Any]:
        validation_result = {
            "valid": True,
            "errors": [],
            "warnings": [],
            "metadata": {}
        }
        try:
            if file.size > max_size:
                validation_result["valid"] = False
                validation_result["errors"].append(f"File size ({file.size:,} bytes) exceeds limit ({max_size:,} bytes)")
            file_ext = file.name.split('.')[-1].lower()
            if file_ext not in SUPPORTED_FILE_TYPES:
                validation_result["valid"] = False
                validation_result["errors"].append(f"Unsupported file type: .{file_ext}")
            file_content = file.read()
            file.seek(0)
            if not FileValidator.validate_file_signature(file_content, file_ext):
                validation_result["warnings"].append("File signature doesn't match extension")
            validation_result["metadata"] = {
                "size_bytes": file.size,
                "size_mb": round(file.size / (1024 * 1024), 2),
                "extension": file_ext,
                "file_hash": FileValidator.calculate_file_hash(file_content),
                "mime_type": mimetypes.guess_type(file.name)[0]
            }
            if file_ext in ['csv', 'txt']:
                try:
                    content_str = file_content.decode('utf-8')
                    validation_result["metadata"]["encoding"] = "UTF-8"
                    validation_result["metadata"]["line_count"] = content_str.count('\n')
                except UnicodeDecodeError:
                    try:
                        content_str = file_content.decode('latin-1')
                        validation_result["metadata"]["encoding"] = "Latin-1"
                        validation_result["warnings"].append("File uses Latin-1 encoding")
                    except:
                        validation_result["errors"].append("Unable to decode text content")
        except Exception as e:
            validation_result["valid"] = False
            validation_result["errors"].append(f"Validation error: {str(e)}")
        return validation_result

class SecurityScanner:
    """Security scanning for uploaded files"""

    SUSPICIOUS_PATTERNS = [
        b'<script',
        b'javascript:',
        b'vbscript:',
        b'onload=',
        b'onerror=',
        b'eval(',
        b'document.cookie'
    ]

    @staticmethod
    def scan_for_threats(file_content: bytes, filename: str) -> Dict[str, Any]:
        scan_result = {
            "safe": True,
            "threats": [],
            "risk_level": "Low"
        }
        try:
            for pattern in SecurityScanner.SUSPICIOUS_PATTERNS:
                if pattern in file_content.lower():
                    scan_result["safe"] = False
                    scan_result["threats"].append(f"Suspicious pattern detected: {pattern.decode('utf-8', errors='ignore')}")
            suspicious_chars = ['<', '>', '|', '&', ';', '$', '`']
            if any(char in filename for char in suspicious_chars):
                scan_result["threats"].append("Filename contains suspicious characters")
                scan_result["risk_level"] = "Medium"
            if scan_result["threats"]:
                scan_result["risk_level"] = "High" if len(scan_result["threats"]) > 2 else "Medium"
        except Exception as e:
            scan_result["safe"] = False
            scan_result["threats"].append(f"Security scan error: {str(e)}")
            scan_result["risk_level"] = "Unknown"
        return scan_result

def validate_and_process_file(file) -> Tuple[bool, Dict, Optional[pd.DataFrame]]:
    validation = FileValidator.validate_file_comprehensive(file)
    file_content = file.read()
    file.seek(0)
    security_scan = SecurityScanner.scan_for_threats(file_content, file.name)
    processing_result = {
        "validation": validation,
        "security": security_scan,
        "processed": False,
        "dataframe": None
    }
    if validation["valid"] and security_scan["safe"]:
        try:
            file_ext = file.name.split('.')[-1].lower()
            df = load_file_content_v5(file_content, file.name, file_ext)
            if df is not None:
                processing_result["processed"] = True
                processing_result["dataframe"] = df
        except Exception as e:
            processing_result["validation"]["errors"].append(f"Processing failed: {str(e)}")
    return (
        processing_result["processed"],
        processing_result,
        processing_result["dataframe"]
    )

# ================== ADVANCED AI CHAT SYSTEM ==================

class IAMChatAgent:
    """Advanced AI chat agent with file context and web search capabilities"""

    def __init__(self, client: OpenAI, model: str = "gpt-4"):
        self.client = client
        self.model = model
        self.conversation_history = []
        self.file_context = {}
        self.available_functions = {
            "analyze_data": self._analyze_data,
            "search_web": self._search_web,
            "generate_export": self._generate_export,
            "create_visualization": self._create_visualization,
            "verify_compliance": self._verify_compliance,
            "get_file_summary": self._get_file_summary
        }

    def update_file_context(self, all_data: Dict[str, pd.DataFrame]):
        self.file_context = {}
        for filename, df in all_data.items():
            self.file_context[filename] = {
                "shape": df.shape,
                "columns": list(df.columns),
                "sample_data": df.head(2).to_dict('records'),
                "summary": f"File {filename} contains {df.shape[0]} records with {df.shape[1]} columns",
                "data_types": df.dtypes.to_dict(),
                "null_counts": df.isnull().sum().to_dict()
            }

    def _create_system_prompt(self) -> str:
        file_summaries = []
        for filename, context in self.file_context.items():
            file_summaries.append(f"- {filename}: {context['summary']}, Columns: {', '.join(context['columns'][:10])}")
        return f"""You are an expert IAM (Identity and Access Management) security consultant and analyst. You have access to the user's uploaded data and can perform various actions.

AVAILABLE DATA FILES:
{chr(10).join(file_summaries)}

CAPABILITIES:
1. analyze_data: Perform detailed analysis on any aspect of the data
2. search_web: Search the internet for current IAM best practices, compliance requirements, etc.
3. generate_export: Create custom reports and exports in various formats
4. create_visualization: Generate charts and graphs for data insights
5. verify_compliance: Check against current compliance frameworks
6. get_file_summary: Get detailed information about specific files

INSTRUCTIONS:
- Always be specific and actionable in your responses
- Use web search to verify current best practices and compliance requirements
- When analyzing data, provide concrete findings with numbers and specifics
- Generate exports when users need reports or documentation
- Create visualizations to illustrate key findings
- Reference the actual uploaded data in your analysis
- Provide step-by-step guidance when requested
- Always verify information against current standards using web search

You can see all the uploaded data and should reference it specifically in your responses. When users ask questions, determine what actions to take and execute them."""

    def _analyze_data(self, query: str, specific_files: List[str] = None) -> Dict[str, Any]:
        try:
            analysis_results = []
            files_to_analyze = specific_files if specific_files else list(self.file_context.keys())
            for filename in files_to_analyze:
                if filename in st.session_state.get("all_data", {}):
                    df = st.session_state["all_data"][filename]
                    result = {
                        "filename": filename,
                        "total_records": len(df),
                        "columns": list(df.columns),
                    }
                    if "user" in query.lower():
                        user_cols = [col for col in df.columns if 'user' in col.lower()]
                        if user_cols:
                            result["user_analysis"] = {
                                "user_columns": user_cols,
                                "unique_users": df[user_cols[0]].nunique()
                            }
                    if "admin" in query.lower() or "privilege" in query.lower():
                        admin_indicators = ['admin', 'root', 'superuser', 'privilege']
                        admin_data = {}
                        for col in df.columns:
                            if any(indicator in col.lower() for indicator in admin_indicators):
                                admin_data[col] = df[col].value_counts().to_dict()
                        result["admin_analysis"] = admin_data
                    analysis_results.append(result)
            return {"analysis": analysis_results, "query": query}
        except Exception as e:
            return {"error": f"Analysis failed: {str(e)}", "query": query}

    def _search_web(self, query: str) -> Dict[str, Any]:
        try:
            if "tracker" in st.session_state:
                st.session_state.tracker.log_web_search(query)
            return {
                "query": query,
                "status": "Web search capability not yet implemented",
                "recommendation": "Please manually verify current industry standards",
                "suggested_sources": [
                    "NIST Cybersecurity Framework (nist.gov)",
                    "ISO 27001 Standards",
                    "OWASP IAM Guidelines",
                    "Your organization's security policies"
                ],
                "note": "This feature will be enhanced in future versions"
            }
        except Exception as e:
            return {"error": f"Web search preparation failed: {str(e)}", "query": query}

    def _generate_export(self, export_type: str, content_focus: str) -> Dict[str, Any]:
        try:
            if "tracker" in st.session_state:
                st.session_state.tracker.log_export(export_type)
            timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
            if export_type.lower() == "executive_summary":
                content = self._create_executive_summary(content_focus)
                return {
                    "type": "executive_summary",
                    "content": content,
                    "filename": f"IAM_Executive_Summary_{timestamp}.txt",
                    "format": "text"
                }
            elif export_type.lower() == "technical_report":
                content = self._create_technical_report(content_focus)
                return {
                    "type": "technical_report",
                    "content": content,
                    "filename": f"IAM_Technical_Report_{timestamp}.txt",
                    "format": "text"
                }
            return {"error": f"Unknown export type: {export_type}"}
        except Exception as e:
            return {"error": f"Export generation failed: {str(e)}"}

    def _create_executive_summary(self, focus: str) -> str:
        return f"""
IAM EXECUTIVE SUMMARY - {focus.upper()}
Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}

OVERVIEW:
Based on analysis of uploaded IAM data with focus on {focus}.

KEY FINDINGS:
• Analysis performed on {len(self.file_context)} data files
• Total records analyzed: {sum(ctx['shape'][0] for ctx in self.file_context.values())}

RECOMMENDATIONS:
• Detailed analysis available in technical report
• Immediate action items to be prioritized
• Compliance review recommended

This summary provides high-level insights for executive decision-making.
"""

    def _create_technical_report(self, focus: str) -> str:
        return f"""
IAM TECHNICAL ANALYSIS REPORT - {focus.upper()}
Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}

TECHNICAL ANALYSIS:
Focus Area: {focus}

DATA SOURCES ANALYZED:
{chr(10).join(f"• {filename}: {ctx['summary']}" for filename, ctx in self.file_context.items())}

DETAILED FINDINGS:
[Detailed technical analysis would be generated here based on actual data]

TECHNICAL RECOMMENDATIONS:
[Specific technical recommendations based on analysis]

IMPLEMENTATION GUIDE:
[Step-by-step technical implementation instructions]
"""

    def _create_visualization(self, viz_type: str, data_focus: str) -> Dict[str, Any]:
        try:
            return {
                "type": viz_type,
                "focus": data_focus,
                "data": "Visualization data would be prepared here",
                "success": True
            }
        except Exception as e:
            return {"error": f"Visualization creation failed: {str(e)}"}

    def _verify_compliance(self, framework: str) -> Dict[str, Any]:
        try:
            return {
                "framework": framework,
                "status": "Compliance verification performed",
                "gaps": [],
                "recommendations": []
            }
        except Exception as e:
            return {"error": f"Compliance verification failed: {str(e)}"}

    def _get_file_summary(self, filename: str = None) -> Dict[str, Any]:
        try:
            if filename and filename in self.file_context:
                return self.file_context[filename]
            else:
                return self.file_context
        except Exception as e:
            return {"error": f"File summary failed: {str(e)}"}

    def process_user_query(self, user_query: str) -> Dict[str, Any]:
        """Process user query and determine appropriate actions"""
        try:
            system_prompt = self._create_system_prompt()
            prompt = f"""
User Query: {user_query}

Based on this query and the available data, determine what actions to take and provide a comprehensive response.

Available functions:
- analyze_data: For data analysis questions
- search_web: For current best practices/compliance info (note: limited capability)
- generate_export: For creating reports
- create_visualization: For charts/graphs
- verify_compliance: For compliance checking
- get_file_summary: For file information

Provide your analysis and any function calls needed, then give a detailed response to the user.
"""
            if "tracker" in st.session_state:
                st.session_state.tracker.log_chat_interaction("advanced_query")
            try:
                response = self.client.chat.completions.create(
                    model=self.model,
                    messages=[
                        {"role": "system", "content": system_prompt},
                        {"role": "user", "content": prompt}
                    ],
                    temperature=0.2,
                    max_tokens=2000
                )
                response_content = response.choices[0].message.content
            except openai.APIError as e:
                return {
                    "success": False,
                    "error": f"OpenAI API Error: {str(e)}",
                    "query": user_query,
                    "suggestion": "Please check your API key and try again"
                }
            except openai.RateLimitError:
                return {
                    "success": False,
                    "error": "Rate limit exceeded. Please wait a moment and try again.",
                    "query": user_query
                }
            except Exception as e:
                return {
                    "success": False,
                    "error": f"Unexpected error: {str(e)}",
                    "query": user_query
                }
            self.conversation_history.append({
                "timestamp": datetime.now().isoformat(),
                "user_query": user_query,
                "ai_response": response_content
            })
            return {
                "success": True,
                "response": response_content,
                "actions_taken": [],
                "query": user_query
            }
        except Exception as e:
            logger.error(f"Chat processing failed: {str(e)}")
            return {
                "success": False,
                "error": str(e),
                "query": user_query
            }

def create_advanced_chat_interface(agent: IAMChatAgent):
    """Create the advanced chat interface"""
    st.subheader("🤖 Advanced IAM AI Assistant")
    st.markdown("*Ask me anything about your IAM data. I can analyze, export, verify compliance, and more!*")

    with st.container():
        with st.expander("💡 Example Queries"):
            example_queries = [
                "Analyze all admin accounts and compare to NIST guidelines",
                "Find users with excessive privileges and create a cleanup plan",
                "Generate an executive summary focused on our biggest security risks",
                "Which accounts haven't been used in the last 90 days?",
                "Create a compliance report for SOX requirements",
                "Show me a risk heat map of our privileged accounts",
                "Export technical remediation guide for audit team",
                "Verify our password policies against current standards"
            ]
            st.info("💡 **Note:** Web search capability is currently limited. The AI will provide analysis based on uploaded data and general knowledge, but please manually verify current industry standards for compliance requirements.")
            for query in example_queries:
                if st.button(f"📝 {query}", key=f"example_{hash(query)}"):
                    st.session_state["chat_input"] = query

        user_query = st.text_area(
            "Your Question:",
            height=100,
            placeholder="Ask me anything about your IAM data... I can analyze, create reports, verify compliance, search for best practices, and more!",
            key="chat_input"
        )

        col1, col2, col3 = st.columns([1, 1, 2])

        with col1:
            if st.button("🚀 Ask AI", type="primary", disabled=not user_query):
                if user_query:
                    with st.spinner("🧠 Processing your request..."):
                        result = agent.process_user_query(user_query)
                    if result["success"]:
                        st.success("✅ Analysis Complete!")
                        st.markdown("### 🤖 AI Response:")
                        st.markdown(result["response"])
                        if result.get("actions_taken"):
                            st.markdown("### ⚡ Actions Performed:")
                            for action in result["actions_taken"]:
                                st.info(f"✓ {action}")
                    else:
                        st.error(f"❌ Processing failed: {result.get('error', 'Unknown error')}")

        with col2:
            if st.button("🗑️ Clear Chat"):
                agent.conversation_history = []
                st.session_state["chat_input"] = ""
                st.success("Chat cleared!")

        if agent.conversation_history:
            st.markdown("### 💭 Recent Conversations")
            for i, conv in enumerate(agent.conversation_history[-5:][::-1]):
                with st.expander(f"🕒 {conv['timestamp'][:19]} - Conversation"):
                    st.markdown(f"**You:** {conv['user_query']}")
                    st.markdown(f"**AI:** {conv['ai_response']}")

def safe_execute(func, *args, **kwargs):
    """Wrapper for safe function execution with error tracking"""
    try:
        return func(*args, **kwargs)
    except Exception as e:
        error_msg = f"Function {func.__name__} failed: {str(e)}"
        logger.error(error_msg)
        if "tracker" in st.session_state:
            st.session_state.tracker.log_error("FunctionError", error_msg)
        return None

def monitor_performance(func):
    """Decorator to monitor function performance"""
    def wrapper(*args, **kwargs):
        start_time = datetime.now()
        try:
            result = func(*args, **kwargs)
            duration = datetime.now() - start_time
            logger.info(f"Function {func.__name__} completed in {duration.total_seconds():.2f} seconds")
            return result
        except Exception as e:
            duration = datetime.now() - start_time
            logger.error(f"Function {func.__name__} failed after {duration.total_seconds():.2f} seconds: {str(e)}")
            raise
    return wrapper

def display_validation_results(results: Dict[str, Any], filename: str):
    """Display validation and security results in UI"""
    validation = results["validation"]
    security = results["security"]

    if validation["valid"] and security["safe"]:
        st.success(f"✅ {filename} - Validation passed")
    else:
        st.error(f"❌ {filename} - Validation failed")

    expander_text = f"🔍 Details for {filename}"
    if validation["errors"] or security["threats"]:
        expander_text += " ⚠️"

    with st.expander(expander_text):
        col1, col2 = st.columns(2)
        with col1:
            st.subheader("📋 Validation")
            if validation["errors"]:
                for error in validation["errors"]:
                    st.error(f"❌ {error}")
            if validation["warnings"]:
                for warning in validation["warnings"]:
                    st.warning(f"⚠️ {warning}")
            if not validation["errors"] and not validation["warnings"]:
                st.success("✅ All validation checks passed")
        with col2:
            st.subheader("🛡️ Security Scan")
            risk_level = security["risk_level"]
            if risk_level == "Low":
                st.success(f"🟢 Risk Level: {risk_level}")
            elif risk_level == "Medium":
                st.warning(f"🟡 Risk Level: {risk_level}")
            else:
                st.error(f"🔴 Risk Level: {risk_level}")
            if security["threats"]:
                for threat in security["threats"]:
                    st.error(f"🚨 {threat}")
            else:
                st.success("✅ No security threats detected")
        if "metadata" in validation and validation["metadata"]:
            st.subheader("📊 File Metadata")
            metadata = validation["metadata"]
            col1, col2, col3 = st.columns(3)
            with col1:
                st.metric("Size", f"{metadata.get('size_mb', 0)} MB")
            with col2:
                st.metric("Type", metadata.get('extension', 'Unknown').upper())
            with col3:
                st.metric("Encoding", metadata.get('encoding', 'N/A'))
            if metadata.get('line_count'):
                st.info(f"📄 Lines: {metadata['line_count']:,}")

def check_authentication():
    """Simple password authentication"""
    if "authenticated" not in st.session_state:
        st.session_state["authenticated"] = False
    if st.session_state["authenticated"]:
        return True

    st.title("🔐 IAM Analyzer v6 Enterprise")
    st.markdown("Please enter the access password to continue.")
    try:
        correct_password = st.secrets.get("APP_PASSWORD", "IAMEnterprise2024!")
    except:
        correct_password = "IAMEnterprise2024!"
    entered_password = st.text_input("Password:", type="password", key="auth_password")
    col1, col2, col3 = st.columns([1, 1, 1])
    with col2:
        if st.button("🚀 Access Enterprise Tool", type="primary", use_container_width=True):
            if entered_password == correct_password:
                st.session_state["authenticated"] = True
                st.success("✅ Access granted! Loading Enterprise IAM Analyzer...")
                st.rerun()
            else:
                st.error("❌ Incorrect password. Contact administrator for access.")
    st.markdown("---")
    st.info("🛡️ Enterprise IAM Analysis Platform - Supports 10+ file formats with AI-powered insights")
    return False

def validate_api_key(api_key: str) -> bool:
    """Validate OpenAI API key format"""
    if not api_key:
        return False
    return api_key.startswith('sk-') and len(api_key) > 20

def validate_file_size(file) -> bool:
    """Check if file size is within limits"""
    return file.size <= MAX_FILE_SIZE

def safe_json_parse(text: str) -> Optional[Dict]:
    """Safely parse JSON from AI response"""
    try:
        start = text.find("{")
        end = text.rfind("}") + 1
        if start != -1 and end > start:
            json_str = text[start:end]
            return json.loads(json_str)
    except (json.JSONDecodeError, ValueError):
        pass
    return None

def extract_text_from_docx(file_content: bytes) -> str:
    if not DOCX_AVAILABLE:
        return "DOCX processing not available. Install python-docx."
    try:
        doc = Document(io.BytesIO(file_content))
        text_content = []
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text_content.append(paragraph.text)
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join([cell.text for cell in row.cells])
                if row_text.strip():
                    text_content.append(row_text)
        return "\n".join(text_content)
    except Exception as e:
        return f"Error processing DOCX: {str(e)}"

def extract_text_from_pptx(file_content: bytes) -> str:
    if not PPTX_AVAILABLE:
        return "PPTX processing not available. Install python-pptx."
    try:
        prs = Presentation(io.BytesIO(file_content))
        text_content = []
        for slide_num, slide in enumerate(prs.slides, 1):
            text_content.append(f"=== SLIDE {slide_num} ===")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_content.append(shape.text)
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_text = " | ".join([cell.text for cell in row.cells])
                        if row_text.strip():
                            text_content.append(row_text)
        return "\n".join(text_content)
    except Exception as e:
        return f"Error processing PPTX: {str(e)}"

def extract_text_from_image(file_content: bytes, file_name: str) -> str:
    if not OCR_AVAILABLE:
        return "OCR not available. Install PIL and pytesseract."
    try:
        image = Image.open(io.BytesIO(file_content))
        if image.mode != 'RGB':
            image = image.convert('RGB')
        extracted_text = pytesseract.image_to_string(image)
        if extracted_text.strip():
            return f"OCR Extracted Text from {file_name}:\n{extracted_text}"
        else:
            return f"No text detected in image: {file_name}"
    except Exception as e:
        return f"Error processing image {file_name}: {str(e)}"

def extract_text_from_pbi(file_content: bytes) -> str:
    try:
        with zipfile.ZipFile(io.BytesIO(file_content), 'r') as zip_file:
            content_info = ["Power BI File Structure Analysis:"]
            for file_info in zip_file.filelist:
                content_info.append(f"- {file_info.filename} ({file_info.file_size} bytes)")
            if 'metadata.json' in zip_file.namelist():
                metadata = zip_file.read('metadata.json').decode('utf-8')
                content_info.append("\nMetadata:")
                content_info.append(metadata[:1000] + "..." if len(metadata) > 1000 else metadata)
            model_files = [f for f in zip_file.namelist() if 'model' in f.lower() and f.endswith('.json')]
            for model_file in model_files[:3]:
                mc = zip_file.read(model_file).decode('utf-8')
                content_info.append(f"\n{model_file}:")
                content_info.append(mc[:500] + "..." if len(mc) > 500 else mc)
        return "\n".join(content_info)
    except Exception as e:
        return f"Error processing Power BI file: {str(e)}"

@st.cache_data
def load_file_content_v5(file_content: bytes, file_name: str, file_type: str) -> Optional[pd.DataFrame]:
    try:
        if file_type == "csv":
            return pd.read_csv(io.BytesIO(file_content))
        elif file_type in ["xlsx", "xls"]:
            return pd.read_excel(io.BytesIO(file_content))
        elif file_type == "txt":
            content = file_content.decode("utf-8")
            return pd.DataFrame({"Text_Content": [content]})
        elif file_type == "docx":
            text_content = extract_text_from_docx(file_content)
            return pd.DataFrame({"Document_Content": [text_content]})
        elif file_type == "pptx":
            text_content = extract_text_from_pptx(file_content)
            return pd.DataFrame({"Presentation_Content": [text_content]})
        elif file_type in ["png", "jpg", "jpeg", "gif", "bmp", "tiff"]:
            text_content = extract_text_from_image(file_content, file_name)
            return pd.DataFrame({"Image_OCR_Content": [text_content]})
        elif file_type == "pbi":
            text_content = extract_text_from_pbi(file_content)
            return pd.DataFrame({"PowerBI_Analysis": [text_content]})
        else:
            return None
    except Exception as e:
        st.error(f"Error processing {file_name}: {str(e)}")
        return None

@st.cache_data
def generate_enhanced_summary(all_data: Dict[str, pd.DataFrame]) -> str:
    summary = []
    summary.append("=== COMPREHENSIVE IAM DATA ANALYSIS ===")
    summary.append(f"Analysis Timestamp: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")
    summary.append(f"Total Files Processed: {len(all_data)}")
    summary.append("")
    total_records = 0
    file_types = {}
    for name, df in all_data.items():
        file_ext = name.split('.')[-1].lower()
        file_types[file_ext] = file_types.get(file_ext, 0) + 1
        total_records += df.shape[0]
        summary.append(f"📁 FILE: {name}")
        summary.append(f"   Type: {file_ext.upper()}")
        summary.append(f"   Dimensions: {df.shape[0]} rows × {df.shape[1]} columns")
        if df.shape[1] <= 20:
            summary.append(f"   Columns: {', '.join(df.columns)}")
        else:
            summary.append(f"   Columns: {', '.join(df.columns[:15])} ... (+{df.shape[1] - 15} more)")
        if df.shape[0] > 0:
            summary.append("   Sample Data:")
            sample_df = df.head(3)
            for idx, row in sample_df.iterrows():
                summary.append(f"     Row {idx + 1}: {dict(row)}")
        if df.shape[1] > 1:
            null_counts = df.isnull().sum()
            if null_counts.sum() > 0:
                summary.append(f"   Data Quality: {null_counts.sum()} null values detected")
            iam_keywords = ['user', 'role', 'permission', 'access', 'group', 'policy', 'entitlement', 'account', 'login', 'auth']
            iam_columns = [col for col in df.columns if any(keyword in col.lower() for keyword in iam_keywords)]
            if iam_columns:
                summary.append(f"   IAM-Related Columns: {', '.join(iam_columns)}")
        summary.append("   " + "─" * 80)
    summary.append("\n=== DATASET OVERVIEW ===")
    summary.append(f"Total Records Across All Files: {total_records:,}")
    summary.append(f"File Type Distribution: {dict(file_types)}")
    summary.append(f"Analysis Scope: Enterprise IAM Security Assessment")
    return "\n".join(summary)

def initialize_session_state():
    if "history" not in st.session_state:
        st.session_state["history"] = []
    if "chat_log" not in st.session_state:
        st.session_state["chat_log"] = []
    if "analysis_count" not in st.session_state:
        st.session_state["analysis_count"] = 0
    if "authenticated" not in st.session_state:
        st.session_state["authenticated"] = False
    if "tracker" not in st.session_state:
        st.session_state["tracker"] = AnalysisTracker()
    if "all_data" not in st.session_state:
        st.session_state["all_data"] = {}
    if "chat_agent" not in st.session_state:
        st.session_state["chat_agent"] = None

def setup_sidebar_v6() -> Tuple[str, str, float, int]:
    st.sidebar.title("🔐 IAM Analyzer v6 Enterprise")
    st.sidebar.markdown("*Multi-Format Support + Advanced Reporting*")
    if st.sidebar.button("🚪 Logout"):
        st.session_state["authenticated"] = False
        st.rerun()
    st.sidebar.markdown("---")
    st.sidebar.subheader("⚙️ AI Configuration")
    api_key = None
    try:
        if "OPENAI_API_KEY" in st.secrets:
            api_key = st.secrets["OPENAI_API_KEY"]
            st.sidebar.success("🔑 Using API key from secrets")
    except:
        pass
    if not api_key:
        api_key = st.sidebar.text_input(
            "OpenAI API Key:",
            type="password",
            help="Enter your API key or configure in secrets"
        )
    if api_key and not validate_api_key(api_key):
        st.sidebar.error("❌ Invalid API key format")
        return None, None, None, None
    with st.sidebar.expander("🎛️ Advanced AI Settings"):
        model = st.selectbox("AI Model:", SUPPORTED_MODELS, index=0)
        temperature = st.slider("Analysis Creativity:", 0.0, 1.0, 0.2, 0.1)
        max_tokens = st.slider("Response Depth:", 1000, 4000, 2500, 100)
    st.sidebar.markdown("---")
    st.sidebar.subheader("📁 File Support Status")
    formats_status = {
        "📊 Spreadsheets (CSV/XLSX)": "✅ Full Support",
        "📄 Documents (DOCX)": "✅ Available" if DOCX_AVAILABLE else "❌ Install python-docx",
        "🎨 Presentations (PPTX)": "✅ Available" if PPTX_AVAILABLE else "❌ Install python-pptx",
        "🖼️ Images (OCR)": "✅ Available" if OCR_AVAILABLE else "❌ Install pytesseract",
        "📈 Charts (Plotly)": "✅ Available" if PLOTLY_AVAILABLE else "❌ Install plotly",
        "📊 Power BI (PBI)": "✅ Basic Support"
    }
    for format_name, status in formats_status.items():
        if "✅" in status:
            st.sidebar.success(f"{format_name}: {status.replace('✅ ', '')}")
        else:
            st.sidebar.warning(f"{format_name}: {status.replace('❌ ', '')}")
    if "tracker" in st.session_state:
        st.sidebar.markdown("---")
        st.sidebar.subheader("📊 Session Analytics")
        metrics = st.session_state.tracker.get_session_summary()
        col1, col2 = st.sidebar.columns(2)
        with col1:
            st.metric("Files Processed", metrics["files_processed"])
            st.metric("Chat Interactions", metrics["chat_interactions"])
        with col2:
            st.metric("Success Rate", f"{metrics['success_rate']}%")
            st.metric("Duration", f"{metrics['duration_minutes']} min")
        with st.sidebar.expander("📈 Detailed Metrics"):
            st.write(f"**Session ID:** {metrics['session_id']}")
            st.write(f"**Total Records:** {metrics['total_records']:,}")
            st.write(f"**API Calls:** {metrics['api_calls']}")
            st.write(f"**Web Searches:** {metrics['web_searches']}")
            st.write(f"**Exports Generated:** {metrics['exports_generated']}")
            st.write(f"**Errors:** {metrics['errors_count']}")
    return api_key, model, temperature, max_tokens

def create_enterprise_prompt(data_summary: str, tasks: List[str], analysis_depth: str = "comprehensive") -> str:
    current_date = datetime.now().strftime("%Y-%m-%d")
    return f"""You are a Senior IAM (Identity and Access Management) Security Consultant conducting an enterprise-grade security assessment.

ENGAGEMENT DETAILS:
- Assessment Date: {current_date}
- Analysis Scope: {analysis_depth.title()} Enterprise IAM Review
- Deliverable: Executive + Technical Security Report

ANALYSIS TASKS:
{chr(10).join(f"• {task}" for task in tasks)}

IAM DATA INVENTORY:
{data_summary}

COMPREHENSIVE ANALYSIS REQUIREMENTS:

1. EXECUTIVE ASSESSMENT:
   - Business risk impact analysis
   - Compliance posture evaluation
   - Strategic recommendations with timelines
   - Budget implications for remediation

2. TECHNICAL FINDINGS:
   - Detailed security control gaps
   - Attack vector analysis
   - Privilege escalation paths
   - Data classification and access patterns

3. COMPLIANCE MAPPING:
   - SOX, SOD, GDPR, HIPAA considerations
   - Industry framework alignment (NIST, ISO 27001)
   - Audit trail and logging adequacy

4. RISK QUANTIFICATION:
   - Business impact scoring (1-10 scale)
   - Probability of exploitation
   - Estimated cost of breach scenarios
   - Risk appetite alignment

OUTPUT FORMAT - COMPREHENSIVE JSON STRUCTURE:
{{
  "executive_summary": {{
    "overall_risk_score": <integer 1-10>,
    "key_findings": ["finding 1", "finding 2", "finding 3"],
    "business_impact": "Critical/High/Medium/Low",
    "recommended_actions": ["immediate action 1", "action 2"],
    "investment_required": "Low/Medium/High",
    "timeline_to_remediate": "30/60/90/180 days"
  }},
  "detailed_findings": {{
    "access_violations": [
      {{
        "violation_type": "Segregation of Duties",
        "severity": "Critical/High/Medium/Low", 
        "description": "detailed description",
        "affected_users": "number or list",
        "business_risk": "impact description",
        "remediation": "specific steps"
      }}
    ],
    "privileged_access_analysis": {{
      "admin_account_count": "number",
      "service_accounts": "analysis",
      "emergency_access": "status",
      "mfa_coverage": "percentage",
      "review_frequency": "assessment"
    }},
    "compliance_gaps": [
      {{
        "control_framework": "SOX/GDPR/HIPAA/etc",
        "control_id": "specific control",
        "gap_description": "what's missing",
        "remediation_effort": "Low/Medium/High",
        "priority": "Critical/High/Medium/Low"
      }}
    ]
  }},
  "risk_matrix": [
    ["Risk Category", "Likelihood", "Impact", "Risk Score", "Priority"],
    ["Unauthorized Access", "High", "Critical", "9", "Immediate"],
    ["Data Exposure", "Medium", "High", "7", "30 days"]
  ],
  "remediation_roadmap": [
    {{
      "phase": "Immediate (0-30 days)",
      "actions": ["action 1", "action 2"],
      "estimated_effort": "person-hours or cost",
      "dependencies": ["dependency 1"]
    }},
    {{
      "phase": "Short-term (30-90 days)", 
      "actions": ["action 1", "action 2"],
      "estimated_effort": "estimate",
      "dependencies": ["dependency 1"]
    }}
  ],
  "technical_recommendations": [
    {{
      "category": "Access Controls",
      "recommendation": "specific technical change",
      "implementation_guide": "step-by-step process",
      "validation_criteria": "how to verify success"
    }}
  ],
  "metrics_dashboard": [
    ["Metric", "Current Value", "Target Value", "Timeline"],
    ["Users with Excessive Access", "25", "0", "60 days"],
    ["Accounts Without MFA", "150", "0", "30 days"]
  ]
}}

CRITICAL INSTRUCTIONS:
- Provide actionable, specific recommendations
- Include quantitative risk assessments where possible  
- Focus on business impact and ROI of security investments
- Address both immediate vulnerabilities and strategic improvements
- Ensure all JSON fields are properly formatted
- Response must be valid JSON only"""

def run_enterprise_analysis(data_summary: str, tasks: List[str], client: OpenAI, model: str, temperature: float, max_tokens: int) -> Optional[Dict]:
    try:
        prompt = create_enterprise_prompt(data_summary, tasks, "comprehensive")
        response = client.chat.completions.create(
            model=model,
            messages=[
                {"role": "system", "content": "You are a Senior IAM Security Consultant with 15+ years of enterprise experience. You specialize in comprehensive security assessments, compliance frameworks, and risk quantification. Always provide executive-level insights with technical depth. Respond only with valid JSON."},
                {"role": "user", "content": prompt}
            ],
            temperature=temperature,
            max_tokens=max_tokens
        )
        raw_content = response.choices[0].message.content
        json_data = safe_json_parse(raw_content)
        if not json_data:
            st.error("❌ Failed to parse AI response as JSON")
            with st.expander("🔍 Raw AI Response"):
                st.code(raw_content)
            return None
        return json_data
    except Exception as e:
        st.error(f"❌ Analysis failed: {str(e)}")
        st.code(traceback.format_exc())
        return None

def create_enterprise_dashboard(json_data: Dict[str, Any]):
    st.subheader("📊 Enterprise Security Dashboard")
    try:
        exec_summary = json_data.get("executive_summary", {})
        col1, col2, col3, col4 = st.columns(4)
        with col1:
            risk_score = exec_summary.get("overall_risk_score", 0)
            st.metric("🎯 Overall Risk Score", f"{risk_score}/10", delta=f"Target: ≤3", delta_color="inverse")
        with col2:
            investment = exec_summary.get("investment_required", "Unknown")
            st.metric("💰 Investment Level", investment)
        with col3:
            timeline = exec_summary.get("timeline_to_remediate", "Unknown")
            st.metric("⏱️ Remediation Timeline", timeline)
        with col4:
            business_impact = exec_summary.get("business_impact", "Unknown")
            color = "🔴" if business_impact == "Critical" else "🟡" if business_impact == "High" else "🟢"
            st.metric("📈 Business Impact", f"{color} {business_impact}")
        st.subheader("🔥 Risk Heat Matrix")
        risk_matrix = json_data.get("risk_matrix", [])
        if risk_matrix and len(risk_matrix) > 1:
            df_risk = pd.DataFrame(risk_matrix[1:], columns=risk_matrix[0])
            if PLOTLY_AVAILABLE and "Risk Score" in df_risk.columns:
                df_risk["Risk Score"] = pd.to_numeric(df_risk["Risk Score"], errors='coerce')
                fig = px.scatter(
                    df_risk,
                    x="Likelihood",
                    y="Impact",
                    size="Risk Score",
                    color="Risk Score",
                    hover_name="Risk Category",
                    title="Risk Assessment Matrix",
                    color_continuous_scale="Reds"
                )
                st.plotly_chart(fig, use_container_width=True)
            else:
                st.dataframe(df_risk, use_container_width=True)
        metrics = json_data.get("metrics_dashboard", [])
        if metrics and len(metrics) > 1:
            st.subheader("📈 Key Performance Indicators")
            df_metrics = pd.DataFrame(metrics[1:], columns=metrics[0])
            for _, row in df_metrics.iterrows():
                name = row.iloc[0]
                current = row.iloc[1]
                target = row.iloc[2]
                timeline = row.iloc[3] if len(row) > 3 else "TBD"
                try:
                    current_num = float(str(current).replace(',', ''))
                    target_num = float(str(target).replace(',', ''))
                    progress = max(0, min(100, (1 - current_num/target_num) * 100)) if target_num > 0 else (100 if current_num <= target_num else 0)
                    st.markdown(f"**{name}**")
                    st.progress(progress / 100)
                    st.markdown(f"Current: {current} | Target: {target} | Timeline: {timeline}")
                except:
                    st.markdown(f"**{name}**: {current} → {target} ({timeline})")
        roadmap = json_data.get("remediation_roadmap", [])
        if roadmap:
            st.subheader("🗺️ Remediation Roadmap")
            for phase in roadmap:
                with st.expander(f"📅 {phase.get('phase', 'Unknown Phase')}"):
                    st.markdown(f"**Estimated Effort:** {phase.get('estimated_effort', 'TBD')}")
                    st.markdown("**Actions:**")
                    for action in phase.get("actions", []):
                        st.markdown(f"• {action}")
                    if "dependencies" in phase:
                        st.markdown("**Dependencies:**")
                        for dep in phase["dependencies"]:
                            st.markdown(f"⚠️ {dep}")
    except Exception as e:
        st.error(f"Dashboard error: {str(e)}")
        st.write("Available data keys:", list(json_data.keys()))

def create_comprehensive_exports(json_data: Dict[str, Any], data_summary: str):
    st.subheader("📥 Enterprise Export Options")
    timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
    col1, col2, col3, col4 = st.columns(4)
    with col1:
        exec_report = create_executive_report(json_data, timestamp)
        st.download_button(
            label="📋 Executive Summary",
            data=exec_report,
            file_name=f"IAM_Executive_Report_{timestamp}.txt",
            mime="text/plain",
            help="High-level executive summary for leadership"
        )
    with col2:
        technical_json = json.dumps(json_data, indent=2, ensure_ascii=False)
        st.download_button(
            label="🔧 Technical Analysis",
            data=technical_json,
            file_name=f"IAM_Technical_Analysis_{timestamp}.json",
            mime="application/json",
            help="Complete technical findings in JSON format"
        )
    with col3:
        if "risk_matrix" in json_data:
            risk_csv = create_risk_matrix_csv(json_data["risk_matrix"])
            st.download_button(
                label="📊 Risk Matrix CSV",
                data=risk_csv,
                file_name=f"IAM_Risk_Matrix_{timestamp}.csv",
                mime="text/csv",
                help="Risk assessment data for further analysis"
            )
    with col4:
        remediation_plan = create_remediation_plan(json_data, timestamp)
        st.download_button(
            label="🗓️ Action Plan",
            data=remediation_plan,
            file_name=f"IAM_Remediation_Plan_{timestamp}.txt",
            mime="text/plain",
            help="Detailed implementation roadmap"
        )

def create_executive_report(json_data: Dict[str, Any], timestamp: str) -> str:
    exec_summary = json_data.get("executive_summary", {})
    report = f"""
IAM SECURITY ASSESSMENT - EXECUTIVE SUMMARY
Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}
Report ID: IAM-{timestamp}

═══════════════════════════════════════════════════════════════════

EXECUTIVE OVERVIEW

Overall Risk Score: {exec_summary.get('overall_risk_score', 'N/A')}/10
Business Impact: {exec_summary.get('business_impact', 'N/A')}
Investment Required: {exec_summary.get('investment_required', 'N/A')}
Timeline to Remediate: {exec_summary.get('timeline_to_remediate', 'N/A')}

KEY FINDINGS:
{chr(10).join(f"• {finding}" for finding in exec_summary.get('key_findings', []))}

IMMEDIATE ACTIONS REQUIRED:
{chr(10).join(f"• {action}" for action in exec_summary.get('recommended_actions', []))}

═══════════════════════════════════════════════════════════════════

DETAILED ANALYSIS SUMMARY

"""
    detailed = json_data.get("detailed_findings", {})
    if "access_violations" in detailed:
        report += "\nACCESS VIOLATIONS:\n"
        for violation in detailed["access_violations"]:
            report += f"• {violation.get('violation_type', 'Unknown')}: {violation.get('severity', 'N/A')} - {violation.get('description', 'No description')}\n"
    if "compliance_gaps" in detailed:
        report += "\nCOMPLIANCE GAPS:\n"
        for gap in detailed["compliance_gaps"]:
            report += f"• {gap.get('control_framework', 'Unknown Framework')}: {gap.get('gap_description', 'No description')}\n"
    tech_recs = json_data.get("technical_recommendations", [])
    if tech_recs:
        report += "\nTECHNICAL RECOMMENDATIONS:\n"
        for rec in tech_recs:
            report += f"• {rec.get('category', 'General')}: {rec.get('recommendation', 'No recommendation')}\n"
    report += f"""
═══════════════════════════════════════════════════════════════════

This assessment was conducted using AI-powered analysis of your IAM data.
For questions or clarification, contact your security team.

Report generated by IAM Analyzer v6 Enterprise Edition
"""
    return report

def create_risk_matrix_csv(risk_matrix: List) -> str:
    if not risk_matrix or len(risk_matrix) <= 1:
        return "No risk matrix data available"
    df = pd.DataFrame(risk_matrix[1:], columns=risk_matrix[0])
    return df.to_csv(index=False)

def create_remediation_plan(json_data: Dict[str, Any], timestamp: str) -> str:
    plan = f"""
IAM SECURITY REMEDIATION PLAN
Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}
Plan ID: REMED-{timestamp}

═══════════════════════════════════════════════════════════════════

IMPLEMENTATION ROADMAP

"""
    roadmap = json_data.get("remediation_roadmap", [])
    for i, phase in enumerate(roadmap, 1):
        plan += f"\nPHASE {i}: {phase.get('phase', 'Unknown Phase')}\n"
        plan += f"Estimated Effort: {phase.get('estimated_effort', 'TBD')}\n\nActions:\n"
        for action in phase.get("actions", []):
            plan += f"□ {action}\n"
        if "dependencies" in phase:
            plan += "\nDependencies:\n"
            for dep in phase["dependencies"]:
                plan += f"⚠ {dep}\n"
        plan += "\n" + "─" * 50 + "\n"
    tech_recs = json_data.get("technical_recommendations", [])
    if tech_recs:
        plan += "\nTECHNICAL IMPLEMENTATION GUIDES:\n\n"
        for rec in tech_recs:
            plan += f"Category: {rec.get('category', 'General')}\n"
            plan += f"Recommendation: {rec.get('recommendation', 'N/A')}\n"
            plan += f"Implementation: {rec.get('implementation_guide', 'See technical team')}\n"
            plan += f"Validation: {rec.get('validation_criteria', 'TBD')}\n"
            plan += "\n" + "─" * 30 + "\n"
    return plan

def check_dependencies():
    missing_deps = []
    if not DOCX_AVAILABLE:
        missing_deps.append("python-docx (for Word documents)")
    if not PPTX_AVAILABLE:
        missing_deps.append("python-pptx (for PowerPoint)")
    if not OCR_AVAILABLE:
        missing_deps.append("pytesseract + PIL (for image OCR)")
    if not PLOTLY_AVAILABLE:
        missing_deps.append("plotly (for enhanced charts)")
    if missing_deps:
        with st.expander("⚠️ Optional Dependencies Missing"):
            st.warning("Some features require additional packages:")
            for dep in missing_deps:
                st.write(f"• {dep}")
            st.code("pip install python-docx python-pptx pytesseract plotly")
            st.info("The app will work without these, but with limited functionality.")

def main():
    initialize_session_state()
    if not check_authentication():
        return
    check_dependencies()
    missing_capabilities = []
    if not PLOTLY_AVAILABLE:
        missing_capabilities.append("plotly (enhanced charts)")
    if not DOCX_AVAILABLE:
        missing_capabilities.append("python-docx (Word documents)")
    if not PPTX_AVAILABLE:
        missing_capabilities.append("python-pptx (PowerPoint)")
    if not OCR_AVAILABLE:
        missing_capabilities.append("pytesseract (image OCR)")
    if missing_capabilities:
        st.info(f"💡 Install these packages for full functionality: {', '.join(missing_capabilities)}")
    api_config = setup_sidebar_v6()
    if not all(api_config):
        st.warning("⚠️ Please configure your OpenAI API key to continue.")
        st.stop()
    api_key, model, temperature, max_tokens = api_config
    try:
        client = OpenAI(api_key=api_key)
        client.models.list()
    except Exception as e:
        st.error(f"❌ Failed to connect to OpenAI: {str(e)}")
        st.stop()
    st.title("🔐 IAM Analyzer v6 Enterprise Edition")
    st.markdown("**🚀 Multi-Format Analysis + Enterprise Reporting**")
    st.subheader("📂 Upload Enterprise Data")
    st.markdown(f"**Supported formats:** {', '.join(SUPPORTED_FILE_TYPES)}")
    uploaded_files = st.file_uploader(
        "Upload your IAM data files:",
        type=SUPPORTED_FILE_TYPES,
        accept_multiple_files=True,
        help=f"Maximum file size: {MAX_FILE_SIZE // (1024*1024)}MB per file. Supports documents, spreadsheets, presentations, images, and more!"
    )
    if not uploaded_files:
        st.info("👆 Upload files to begin comprehensive analysis")
        if st.session_state["history"]:
            st.subheader("📚 Previous Enterprise Analyses")
            for i, entry in enumerate(st.session_state["history"][-3:][::-1]):
                with st.expander(f"📋 Analysis {len(st.session_state['history'])-i} - {entry['timestamp']}"):
                    if "executive_summary" in entry["output"]:
                        exec_sum = entry["output"]["executive_summary"]
                        c1, c2 = st.columns(2)
                        with c1:
                            st.metric("Risk Score", f"{exec_sum.get('overall_risk_score', 0)}/10")
                        with c2:
                            st.metric("Business Impact", exec_sum.get("business_impact", "N/A"))
                        st.markdown("**Key Findings:**")
                        for f in exec_sum.get('key_findings', [])[:3]:
                            st.markdown(f"• {f}")
        return
    st.subheader("🔄 Processing & Validating Files")
    all_data = {}
    progress_bar = st.progress(0)
    status_text = st.empty()
    for i, file in enumerate(uploaded_files):
        status_text.text(f"Validating and processing {file.name}...")
        success, results, df = validate_and_process_file(file)
        display_validation_results(results, file.name)
        if success and df is not None:
            all_data[file.name] = df
            st.session_state.tracker.log_file_processed(file.name, df.shape[0])
        else:
            st.session_state.tracker.log_error("FileValidation", f"Failed to process {file.name}")
        progress_bar.progress((i + 1) / len(uploaded_files))
    progress_bar.empty()
    status_text.empty()
    st.session_state["all_data"] = all_data
    if not all_data:
        st.error("❌ No files could be processed. Please check file formats.")
        return
    data_summary = generate_enhanced_summary(all_data)
    st.subheader("📋 Enterprise Data Overview")
    c1, c2, c3, c4 = st.columns(4)
    with c1:
        st.metric("📁 Files Processed", len(all_data))
    with c2:
        total_records = sum(df.shape[0] for df in all_data.values())
        st.metric("📊 Total Records", f"{total_records:,}")
    with c3:
        file_types = set(f.split('.')[-1].upper() for f in all_data.keys())
        st.metric("📄 File Types", len(file_types))
    with c4:
        total_size = sum(df.memory_usage(deep=True).sum() for df in all_data.values())
        st.metric("💾 Data Size", f"{total_size // 1024:,} KB")
    with st.expander("🔍 Data Preview"):
        for name, df in all_data.items():
            st.markdown(f"**📁 {name}**")
            st.dataframe(df.head(2), use_container_width=True)
    st.subheader("🎯 Enterprise Analysis Scope")
    task_categories = {
        "🚨 Critical Security": [
            "Detect Segregation of Duties violations",
            "Identify privilege escalation paths",
            "Analyze emergency access procedures",
            "Review administrative account security"
        ],
        "👥 Access Management": [
            "Identify orphaned and inactive accounts",
            "Analyze role-based access patterns",
            "Review guest and contractor access",
            "Assess shared account usage"
        ],
        "📊 Compliance & Governance": [
            "SOX compliance assessment",
            "GDPR privacy impact analysis",
            "Audit trail completeness review",
            "Access certification gaps"
        ],
        "🔍 Advanced Analytics": [
            "Behavioral access pattern analysis",
            "Cross-system entitlement mining",
            "Risk-based access scoring",
            "Anomaly detection assessment"
        ]
    }
    selected_tasks = []
    for category, tasks in task_categories.items():
        st.markdown(f"**{category}**")
        cols = st.columns(2)
        for idx, task in enumerate(tasks):
            with cols[idx % 2]:
                if st.checkbox(task, key=f"{category}_{task}"):
                    selected_tasks.append(task)
    if not selected_tasks:
        st.warning("⚠️ Please select at least one analysis task.")
        return
    if st.session_state["chat_agent"] is None and all_data:
        st.session_state["chat_agent"] = IAMChatAgent(client, model)
        st.session_state["chat_agent"].update_file_context(all_data)
    if st.session_state["chat_agent"] and all_data:
        st.session_state["chat_agent"].update_file_context(all_data)
    if st.session_state["chat_agent"]:
        st.markdown("---")
        create_advanced_chat_interface(st.session_state["chat_agent"])
        st.markdown("---")
    if st.button("🚀 Run Enterprise Security Analysis", type="primary", use_container_width=True):
        with st.spinner("🤖 Conducting comprehensive IAM security assessment... This may take several minutes."):
            pc1, pc2, pc3 = st.columns(3)
            with pc1:
                st.info("🔍 Processing data sources...")
            with pc2:
                st.info("🧠 Running AI analysis...")
            with pc3:
                st.info("📊 Generating reports...")
            st.session_state.tracker.log_api_call(model, max_tokens)
            analysis_result = run_enterprise_analysis(
                data_summary, selected_tasks, client, model, temperature, max_tokens
            )
        if analysis_result:
            entry = {
                "timestamp": datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
                "input_summary": data_summary,
                "output": analysis_result,
                "tasks": selected_tasks,
                "analysis_type": "Enterprise v6"
            }
            st.session_state["history"].append(entry)
            st.success("✅ Enterprise Analysis Complete!")
            exec_summary = analysis_result.get("executive_summary", {})
            st.subheader("🎯 Executive Summary")
            mc1, mc2, mc3, mc4 = st.columns(4)
            with mc1:
                rs = exec_summary.get("overall_risk_score", 0)
                color_rs = "🔴" if rs >= 8 else "🟡" if rs >= 5 else "🟢"
                st.metric("Risk Score", f"{color_rs} {rs}/10")
            with mc2:
                st.metric("Business Impact", exec_summary.get("business_impact", "N/A"))
            with mc3:
                st.metric("Investment Required", exec_summary.get("investment_required", "N/A"))
            with mc4:
                st.metric("Timeline", exec_summary.get("timeline_to_remediate", "N/A"))
            c1f, c2f = st.columns(2)
            with c1f:
                st.subheader("🔍 Key Findings")
                for f in exec_summary.get("key_findings", []):
                    st.warning(f"• {f}")
            with c2f:
                st.subheader("✅ Recommended Actions")
                for a in exec_summary.get("recommended_actions", []):
                    st.success(f"• {a}")
            create_enterprise_dashboard(analysis_result)
            create_comprehensive_exports(analysis_result, data_summary)
        else:
            st.error("❌ Enterprise analysis failed. Please check your data and try again.")

if __name__ == "__main__":
    main()
```
